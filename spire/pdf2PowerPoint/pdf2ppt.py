from spire.pdf.common import *
from spire.pdf import *
import os
from pptx import Presentation
from pptx.util import Inches, Pt

input_pdf = "Sample.pdf"
output_ppt = "output.pptx"

def extract_text_from_pdf(pdf_path):
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")

    doc = PdfDocument()
    doc.LoadFromFile(pdf_path)

    all_pages_text = []
    for i in range(doc.Pages.Count):
        page = doc.Pages[i]
        extractor = PdfTextExtractor(page)
        options = PdfTextExtractOptions()
        text = extractor.ExtractText(options)
        all_pages_text.append(text)

    doc.Close()
    return all_pages_text

def create_ppt_file(text_pages, filename):
    prs = Presentation()
    for i, page_text in enumerate(text_pages, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = f"Page {i}\n\n{page_text}"
        p.font.size = Pt(14)

    prs.save(filename)
    print(f"PPT file saved: {filename}")

def main():
    try:
        pages_text = extract_text_from_pdf(input_pdf)
        create_ppt_file(pages_text, output_ppt)
        print("PowerPoint presentation created successfully!")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    main()
